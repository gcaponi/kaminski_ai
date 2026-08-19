#!/usr/bin/env python3
"""Ingest files dropped in raw/inbox/ (Kaminski does not run this).

Supported: .txt .md .pdf .pptx .docx
Writes a wiki note in raw/docs/ and classifies it.

Usage:
    python scripts/ingest_inbox.py
"""

from __future__ import annotations

import hashlib
import re
import shutil
import subprocess
import sys
from datetime import date
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
INBOX = ROOT / "raw" / "inbox"
DOCS = ROOT / "raw" / "docs"
DONE = INBOX / "processados"
ASK = INBOX / "perguntas"
sys.path.insert(0, str(ROOT / "scripts"))
from themes import kaminski_questions, needs_kaminski, pick_themes  # noqa: E402

SKIP = {".gitkeep", "LEIA-ME.md", "README.md"}


def slugify(name: str) -> str:
    s = Path(name).stem.lower()
    s = re.sub(r"[^a-z0-9]+", "-", s).strip("-")
    return s[:80] or "documento"


def extract_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def extract_pdf(path: Path) -> str:
    try:
        import pymupdf  # type: ignore
    except ImportError:
        return ""
    doc = pymupdf.open(path)
    return "\n".join(page.get_text() for page in doc)


def extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return ""
    prs = Presentation(str(path))
    chunks: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                chunks.append(shape.text)
    return "\n".join(chunks)


def extract_docx(path: Path) -> str:
    try:
        import docx  # type: ignore
    except ImportError:
        return ""
    d = docx.Document(str(path))
    return "\n".join(p.text for p in d.paragraphs)


EXTRACTORS = {
    ".txt": extract_txt,
    ".md": extract_txt,
    ".pdf": extract_pdf,
    ".pptx": extract_pptx,
    ".docx": extract_docx,
}


def write_doc(slug: str, title: str, source_name: str, body: str, themes: list[str]) -> Path:
    DOCS.mkdir(parents=True, exist_ok=True)
    digest = hashlib.sha256(body.encode("utf-8")).hexdigest()
    links = " · ".join(f"[[{t}]]" for t in themes) + " · [[gabriel-kaminski]]"
    note = "\n".join(
        [
            "---",
            f'title: "{title.replace(chr(34), chr(39))}"',
            f"created: {date.today().isoformat()}",
            f"updated: {date.today().isoformat()}",
            "type: source",
            f"source_file: {source_name}",
            f"language: pt-BR",
            f"temas: [{', '.join(themes)}]",
            f"sha256: {digest}",
            "---",
            "",
            f"# {title}",
            "",
            f"- Origem: `{source_name}`",
            f"- Temas: {links}",
            "",
            "## Texto extraído",
            "",
            body.strip() + "\n",
        ]
    )
    dest = DOCS / f"{slug}.md"
    dest.write_text(note, encoding="utf-8")
    return dest


def main() -> int:
    INBOX.mkdir(parents=True, exist_ok=True)
    DONE.mkdir(parents=True, exist_ok=True)
    files = [
        p
        for p in sorted(INBOX.iterdir())
        if p.is_file() and p.name not in SKIP and p.suffix.lower() in EXTRACTORS
    ]
    if not files:
        print("inbox empty")
        return 0
    n_ok = 0
    n_ask = 0
    ASK.mkdir(parents=True, exist_ok=True)
    for path in files:
        ext = path.suffix.lower()
        text = EXTRACTORS[ext](path)
        if not text.strip():
            print(f"FAIL {path.name}  empty extract")
            continue
        slug = slugify(path.name)
        themes = pick_themes(path.stem, text)
        if needs_kaminski(themes):
            qpath = ASK / f"{date.today().isoformat()}-{slug}.md"
            qpath.write_text(
                "---\n"
                f"status: aguardando-kaminski\n"
                f"source: {path.name}\n"
                f"created: {date.today().isoformat()}\n"
                "---\n\n"
                + kaminski_questions(path.stem, text),
                encoding="utf-8",
            )
            print("ASK_KAMINSKI")
            print(kaminski_questions(path.stem, text))
            print(f"WROTE {qpath.relative_to(ROOT)}")
            n_ask += 1
            continue
        write_doc(slug, path.stem, path.name, text, themes)
        shutil.move(str(path), str(DONE / path.name))
        print(f"OK   {path.name} → raw/docs/{slug}.md  {themes}")
        n_ok += 1
    if n_ok:
        subprocess.run([sys.executable, str(ROOT / "scripts" / "classify.py")], check=False)
    print(f"done ok={n_ok} ask_kaminski={n_ask} total={len(files)}")
    return 2 if n_ask else 0


if __name__ == "__main__":
    raise SystemExit(main())
